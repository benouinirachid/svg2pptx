"""PowerPoint shape creation utilities."""

import math
from typing import Optional

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes, GroupShapes
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from svg2pptx.parser.styles import (
    Style,
    Gradient,
    LinearGradient,
    RadialGradient,
)
from svg2pptx.parser.shapes import (
    ParsedShape,
    RectShape,
    CircleShape,
    EllipseShape,
    LineShape,
    PolygonShape,
    PolylineShape,
)
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.geometry.transforms import Transform


def create_shape(
    shapes: SlideShapes,
    parsed_shape: ParsedShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint shape from a parsed SVG shape.

    Args:
        shapes: SlideShapes or GroupShapes collection to add shape to.
        parsed_shape: Parsed shape data.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    if isinstance(parsed_shape, RectShape):
        return create_rectangle(shapes, parsed_shape, offset_x, offset_y, scale)
    elif isinstance(parsed_shape, (CircleShape, EllipseShape)):
        return create_oval(shapes, parsed_shape, offset_x, offset_y, scale)
    elif isinstance(parsed_shape, LineShape):
        return create_line(shapes, parsed_shape, offset_x, offset_y, scale)
    elif isinstance(parsed_shape, (PolygonShape, PolylineShape)):
        from svg2pptx.pptx_writer.freeform import create_freeform

        return create_freeform(shapes, parsed_shape, offset_x, offset_y, scale)
    return None


def create_rectangle(
    shapes: SlideShapes,
    rect: RectShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> BaseShape:
    """Create a PowerPoint rectangle shape."""
    # Apply transform to get actual position
    x, y = rect.transform.apply(rect.x, rect.y)
    
    # Convert to EMU with scale
    left = offset_x + px_to_emu(x * scale)
    top = offset_y + px_to_emu(y * scale)
    width = px_to_emu(rect.width * scale)
    height = px_to_emu(rect.height * scale)

    # Choose shape type based on corner radius
    if rect.rx > 0 or rect.ry > 0:
        shape = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        # Note: python-pptx doesn't easily support setting corner radius
    else:
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    apply_style(shape, rect.style, rect.transform)
    return shape


def create_oval(
    shapes: SlideShapes,
    oval: CircleShape | EllipseShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> BaseShape:
    """Create a PowerPoint oval (ellipse/circle) shape."""
    if isinstance(oval, CircleShape):
        cx, cy = oval.transform.apply(oval.cx, oval.cy)
        rx = ry = oval.r
    else:
        cx, cy = oval.transform.apply(oval.cx, oval.cy)
        rx, ry = oval.rx, oval.ry

    # Convert center + radius to left, top, width, height
    left = offset_x + px_to_emu((cx - rx) * scale)
    top = offset_y + px_to_emu((cy - ry) * scale)
    width = px_to_emu(2 * rx * scale)
    height = px_to_emu(2 * ry * scale)

    shape = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    apply_style(shape, oval.style, oval.transform)
    return shape


def create_line(
    shapes: SlideShapes,
    line: LineShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> BaseShape:
    """Create a PowerPoint line connector."""
    from pptx.enum.shapes import MSO_CONNECTOR

    # Apply transform
    x1, y1 = line.transform.apply(line.x1, line.y1)
    x2, y2 = line.transform.apply(line.x2, line.y2)

    # Convert to EMU
    start_x = offset_x + px_to_emu(x1 * scale)
    start_y = offset_y + px_to_emu(y1 * scale)
    end_x = offset_x + px_to_emu(x2 * scale)
    end_y = offset_y + px_to_emu(y2 * scale)

    connector = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )

    # Disable shadow on line
    try:
        connector.shadow.inherit = False
        if hasattr(connector.shadow, 'visible'):
            connector.shadow.visible = False
    except (AttributeError, NotImplementedError):
        pass

    # Apply stroke style to line
    if line.style.gradient_stroke is not None:
        _apply_gradient_fill(
            connector.line.fill, line.style.gradient_stroke, line.transform
        )
    elif line.style.stroke != "none":
        try:
            color = parse_hex_color(line.style.stroke)
            connector.line.color.rgb = color
        except ValueError:
            pass

    connector.line.width = Emu(px_to_emu(line.style.stroke_width))

    return connector


_DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _apply_gradient_fill(
    fill, gradient: Gradient, transform: Optional[Transform] = None
) -> None:
    """Apply a linear or radial gradient to a python-pptx FillFormat.

    Works for both shape fills (`shape.fill`) and line fills (`shape.line.fill`),
    since both expose the same FillFormat interface.
    """
    if not gradient.stops:
        fill.background()
        return

    fill.gradient()
    _write_gradient_stops(fill, gradient.stops)

    if isinstance(gradient, RadialGradient):
        _make_radial_geometry(fill, gradient)
    else:
        _apply_linear_angle(fill, gradient, transform)


def _write_gradient_stops(fill, stops) -> None:
    """Replace the default gradient stops with the parsed SVG stops."""
    import lxml.etree as etree

    gsLst = fill.gradient_stops._gsLst

    # python-pptx initialises gradFill with exactly 2 stops; clear them first.
    for child in list(gsLst):
        gsLst.remove(child)

    for stop in stops:
        pos = max(0, min(100000, int(round(stop.offset * 100000))))
        hex_color = stop.color.lstrip("#").upper()
        if len(hex_color) != 6:
            continue
        gs = etree.SubElement(gsLst, f"{{{_DML_NS}}}gs", pos=str(pos))
        srgb = etree.SubElement(gs, f"{{{_DML_NS}}}srgbClr", val=hex_color)
        if stop.opacity < 1.0:
            alpha_val = max(0, min(100000, int(round(stop.opacity * 100000))))
            etree.SubElement(srgb, f"{{{_DML_NS}}}alpha", val=str(alpha_val))


def _apply_linear_angle(
    fill, gradient: LinearGradient, transform: Optional[Transform]
) -> None:
    """Write the <a:lin> direction element for a linear gradient.

    Written directly via lxml rather than python-pptx's ``gradient_angle``
    setter, because ``line.fill.gradient()`` does not create an <a:lin> for the
    setter to update (only ``shape.fill.gradient()`` does).
    """
    import lxml.etree as etree

    # SVG direction (dx,dy, y-down) → angle CCW from East in screen space.
    dx = gradient.x2 - gradient.x1
    dy = gradient.y2 - gradient.y1

    # For userSpaceOnUse gradients the direction vector is in the shape's local
    # coordinate system. Apply the linear part of the cumulative shape transform
    # to map it to screen space so the angle is correct for rotated shapes.
    if transform is not None and gradient.gradient_units == "userSpaceOnUse":
        dx, dy = (
            transform.a * dx + transform.c * dy,
            transform.b * dx + transform.d * dy,
        )

    if dx == 0.0 and dy == 0.0:
        angle_deg = 0.0
    else:
        angle_deg = math.degrees(math.atan2(-dy, dx)) % 360.0

    gsLst = fill.gradient_stops._gsLst
    gradFill = gsLst.getparent()
    for el in gradFill.findall(f"{{{_DML_NS}}}lin"):
        gradFill.remove(el)

    # OOXML <a:lin ang> is measured clockwise in 60000ths of a degree.
    ooxml_ang = int(round((360.0 - angle_deg) * 60000))
    etree.SubElement(
        gradFill, f"{{{_DML_NS}}}lin", ang=str(ooxml_ang), scaled="0"
    )


def _make_radial_geometry(fill, gradient: RadialGradient) -> None:
    """Convert the gradient into an OOXML radial (`<a:path path="circle">`).

    OOXML radial gradients fill inward from the shape's bounding rectangle to a
    focus rectangle, so this is a centered (optionally focus-shifted) circular
    approximation of the SVG radial — cx/cy/fx/fy are honoured for
    objectBoundingBox units, but r and elliptical/non-uniform radii are not.
    """
    import lxml.etree as etree

    gsLst = fill.gradient_stops._gsLst
    gradFill = gsLst.getparent()

    # python-pptx adds an <a:lin>; replace it (and any prior path) with a path.
    for tag in ("lin", "path"):
        for el in gradFill.findall(f"{{{_DML_NS}}}{tag}"):
            gradFill.remove(el)

    path_el = etree.SubElement(gradFill, f"{{{_DML_NS}}}path", path="circle")

    # Focus point as a fraction of the bounding box. For objectBoundingBox the
    # focal/center coords are already in [0,1]; for userSpaceOnUse we lack the
    # shape bbox here, so fall back to a centered focus.
    if gradient.gradient_units == "objectBoundingBox":
        fx = gradient.fx if gradient.fx is not None else gradient.cx
        fy = gradient.fy if gradient.fy is not None else gradient.cy
    else:
        fx = fy = 0.5

    def pct(v: float) -> str:
        return str(max(0, min(100000, int(round(v * 100000)))))

    # fillToRect insets collapse the fill origin to the focus point.
    etree.SubElement(
        path_el,
        f"{{{_DML_NS}}}fillToRect",
        l=pct(fx),
        t=pct(fy),
        r=pct(1.0 - fx),
        b=pct(1.0 - fy),
    )


def apply_style(shape: BaseShape, style: Style, transform: Optional[Transform] = None, disable_shadow: bool = True) -> None:
    """
    Apply SVG style to a PowerPoint shape.

    Args:
        shape: PowerPoint shape to style.
        style: Parsed SVG style.
        transform: Cumulative shape transform, used to rotate gradient direction.
        disable_shadow: Whether to disable shadow on the shape. Defaults to True.
    """
    # Disable shadow if requested
    if disable_shadow:
        try:
            shape.shadow.inherit = False
            # Setting shadow to no shadow by making it transparent
            if hasattr(shape.shadow, 'visible'):
                shape.shadow.visible = False
        except (AttributeError, NotImplementedError):
            # Some shapes may not support shadow property
            pass

    # Apply fill
    fill = shape.fill
    if style.gradient_fill is not None:
        _apply_gradient_fill(fill, style.gradient_fill, transform)
    elif style.fill == "none":
        fill.background()
    else:
        try:
            color = parse_hex_color(style.fill)
            fill.solid()
            fill.fore_color.rgb = color
        except ValueError:
            fill.background()

    # Apply stroke
    line = shape.line
    if style.gradient_stroke is not None:
        _apply_gradient_fill(line.fill, style.gradient_stroke, transform)
        line.width = Emu(px_to_emu(style.stroke_width))
    elif style.stroke == "none":
        line.fill.background()  # No stroke
    else:
        try:
            color = parse_hex_color(style.stroke)
            line.color.rgb = color
            line.width = Emu(px_to_emu(style.stroke_width))
        except ValueError:
            line.fill.background()


def parse_hex_color(hex_color: str) -> RGBColor:
    """
    Parse a hex color string to RGBColor.

    Args:
        hex_color: Color in format "#RRGGBB" or "#RGB".

    Returns:
        RGBColor object.

    Raises:
        ValueError: If color format is invalid.
    """
    if not hex_color or hex_color == "none":
        raise ValueError("Invalid color: none")

    color = hex_color.strip().lstrip("#")

    if len(color) == 3:
        color = "".join(c * 2 for c in color)

    if len(color) != 6:
        raise ValueError(f"Invalid hex color: {hex_color}")

    try:
        r = int(color[0:2], 16)
        g = int(color[2:4], 16)
        b = int(color[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        raise ValueError(f"Invalid hex color: {hex_color}")
