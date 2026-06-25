"""PowerPoint freeform shape creation for polygons, polylines, and paths."""

import math
from typing import Optional, Union

from lxml import etree
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes, GroupShapes
from pptx.util import Emu

from svg2pptx.parser.shapes import PolygonShape, PolylineShape
from svg2pptx.parser.paths import PathShape
from svg2pptx.parser.segments import (
    PathSegment,
    MoveToSeg,
    LineToSeg,
    CubicBezierToSeg,
    QuadBezierToSeg,
    ArcSeg,
    CloseSeg,
    apply_transform,
)
from svg2pptx.parser.styles import Style
from svg2pptx.geometry.units import px_to_emu, EMU_PER_PX
from svg2pptx.geometry.curves import svg_arc_to_center, svg_arc_to_lines
from svg2pptx.pptx_writer.shapes import apply_style


def create_freeform(
    shapes: SlideShapes,
    parsed_shape: Union[PolygonShape, PolylineShape, PathShape],
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint freeform shape from polygon, polyline, or path.

    Args:
        shapes: SlideShapes collection to add shape to.
        parsed_shape: Parsed polygon, polyline, or path shape.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created freeform shape or None.
    """
    if isinstance(parsed_shape, PathShape):
        return create_freeform_from_path(
            shapes, parsed_shape, offset_x, offset_y, scale
        )
    else:
        return create_freeform_from_points(
            shapes, parsed_shape, offset_x, offset_y, scale
        )


def create_freeform_from_points(
    shapes: SlideShapes,
    parsed_shape: Union[PolygonShape, PolylineShape],
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a freeform shape from polygon or polyline points.

    Args:
        shapes: SlideShapes collection.
        parsed_shape: Parsed polygon or polyline.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    points = parsed_shape.points
    if len(points) < 2:
        return None

    transformed_points = parsed_shape.transform.apply_to_points(points)

    emu_points = [
        (px_to_emu(x * scale), px_to_emu(y * scale))
        for x, y in transformed_points
    ]

    is_closed = isinstance(parsed_shape, PolygonShape)

    first_x, first_y = emu_points[0]
    builder = shapes.build_freeform(first_x, first_y)
    builder.add_line_segments(emu_points[1:], close=is_closed)

    shape = builder.convert_to_shape(offset_x, offset_y)
    apply_style(shape, parsed_shape.style)
    return shape


def create_freeform_from_path(
    shapes: SlideShapes,
    path: PathShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a freeform shape from an SVG path.

    Bezier curves and axis-aligned arcs are emitted as native OOXML elements
    (<a:cubicBezTo>, <a:quadBezTo>, <a:arcTo>) so PowerPoint renders them at
    full quality.  Arcs with a non-zero x-axis rotation fall back to line
    segment approximation.

    Args:
        shapes: SlideShapes collection.
        path: Parsed path shape.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    if not path.segments:
        return None

    # Apply SVG transform, then convert pixel coords to EMU.
    transformed = apply_transform(path.transform, path.segments)
    emu_segs = _segments_to_emu(transformed, scale)

    # Bounding box over all key points (endpoints + control points + arc samples).
    min_x, min_y, max_x, max_y = _bounding_box(emu_segs)
    if min_x is None:
        return None

    shape_w = max(1, int(max_x - min_x))
    shape_h = max(1, int(max_y - min_y))
    shape_x = int(offset_x + min_x)
    shape_y = int(offset_y + min_y)

    # Create the freeform <p:sp> element and get its path list.
    spTree = shapes._spTree  # pyright: ignore[reportPrivateUsage]
    sp = spTree.add_freeform_sp(shape_x, shape_y, shape_w, shape_h)
    path_el = sp.add_path(w=Emu(shape_w), h=Emu(shape_h))

    # Write segment elements, tracking pen position in shape-relative EMU.
    pen_x = pen_y = 0.0
    subpath_start_x = subpath_start_y = 0.0

    for seg in emu_segs:
        if isinstance(seg, MoveToSeg):
            rx, ry = seg.x - min_x, seg.y - min_y
            _emit_move_to(path_el, rx, ry)
            pen_x, pen_y = rx, ry
            subpath_start_x, subpath_start_y = rx, ry

        elif isinstance(seg, LineToSeg):
            rx, ry = seg.x - min_x, seg.y - min_y
            _emit_line_to(path_el, rx, ry)
            pen_x, pen_y = rx, ry

        elif isinstance(seg, CubicBezierToSeg):
            rcp1 = (seg.cp1[0] - min_x, seg.cp1[1] - min_y)
            rcp2 = (seg.cp2[0] - min_x, seg.cp2[1] - min_y)
            rend = (seg.end[0] - min_x, seg.end[1] - min_y)
            _emit_cubic_bez_to(path_el, *rcp1, *rcp2, *rend)
            pen_x, pen_y = rend

        elif isinstance(seg, QuadBezierToSeg):
            rcp = (seg.cp[0] - min_x, seg.cp[1] - min_y)
            rend = (seg.end[0] - min_x, seg.end[1] - min_y)
            _emit_quad_bez_to(path_el, *rcp, *rend)
            pen_x, pen_y = rend

        elif isinstance(seg, ArcSeg):
            # Absolute start position of the arc (needed for center computation).
            abs_start_x = pen_x + min_x
            abs_start_y = pen_y + min_y
            native = _try_emit_arc_to(path_el, abs_start_x, abs_start_y, seg)
            if not native:
                pts = svg_arc_to_lines(
                    x1=abs_start_x, y1=abs_start_y,
                    rx=seg.rx, ry=seg.ry,
                    x_axis_rotation=seg.x_rotation,
                    large_arc_flag=seg.large_arc,
                    sweep_flag=seg.sweep,
                    x2=seg.x, y2=seg.y,
                    tolerance=float(EMU_PER_PX),  # 1-pixel accuracy in EMU
                )
                for lx, ly in pts:
                    _emit_line_to(path_el, lx - min_x, ly - min_y)
            pen_x, pen_y = seg.x - min_x, seg.y - min_y

        elif isinstance(seg, CloseSeg):
            _emit_close(path_el)
            pen_x, pen_y = subpath_start_x, subpath_start_y

    shape = shapes._shape_factory(sp)  # pyright: ignore[reportPrivateUsage]
    apply_style(shape, path.style)
    return shape


# ---------------------------------------------------------------------------
# OOXML element emitters
# ---------------------------------------------------------------------------

def _pt_el(parent, x: float, y: float) -> None:
    pt = etree.SubElement(parent, qn("a:pt"))
    pt.set("x", str(int(round(x))))
    pt.set("y", str(int(round(y))))


def _emit_move_to(path_el, x: float, y: float) -> None:
    el = etree.SubElement(path_el, qn("a:moveTo"))
    _pt_el(el, x, y)


def _emit_line_to(path_el, x: float, y: float) -> None:
    el = etree.SubElement(path_el, qn("a:lnTo"))
    _pt_el(el, x, y)


def _emit_cubic_bez_to(
    path_el,
    cp1x: float, cp1y: float,
    cp2x: float, cp2y: float,
    ex: float, ey: float,
) -> None:
    el = etree.SubElement(path_el, qn("a:cubicBezTo"))
    _pt_el(el, cp1x, cp1y)
    _pt_el(el, cp2x, cp2y)
    _pt_el(el, ex, ey)


def _emit_quad_bez_to(
    path_el, cpx: float, cpy: float, ex: float, ey: float
) -> None:
    el = etree.SubElement(path_el, qn("a:quadBezTo"))
    _pt_el(el, cpx, cpy)
    _pt_el(el, ex, ey)


def _emit_arc_to(
    path_el, wR: float, hR: float, stAng: float, swAng: float
) -> None:
    el = etree.SubElement(path_el, qn("a:arcTo"))
    el.set("wR", str(int(round(wR))))
    el.set("hR", str(int(round(hR))))
    el.set("stAng", str(int(round(stAng))))
    el.set("swAng", str(int(round(swAng))))


def _emit_close(path_el) -> None:
    etree.SubElement(path_el, qn("a:close"))


def _try_emit_arc_to(
    path_el,
    start_x: float,
    start_y: float,
    seg: ArcSeg,
) -> bool:
    """
    Emit <a:arcTo> for seg if possible; return False when fallback is needed.

    OOXML <a:arcTo> only supports axis-aligned ellipses, so arcs with a
    non-zero x-axis rotation fall back to line approximation.
    """
    if seg.rx <= 0 or seg.ry <= 0:
        return False

    x_rot = seg.x_rotation % 360
    if x_rot > 180:
        x_rot -= 360
    if abs(x_rot) > 0.01:
        return False

    cx, cy, rx, ry, theta1, dtheta = svg_arc_to_center(
        x1=start_x, y1=start_y,
        rx=seg.rx, ry=seg.ry,
        phi=0.0,
        large_arc=seg.large_arc,
        sweep=seg.sweep,
        x2=seg.x, y2=seg.y,
    )

    # svg_arc_to_center returns (0,0,0,0,0,0) for degenerate arcs.
    if rx <= 0 or ry <= 0:
        return False

    # Angles in 60000ths of a degree (OOXML unit).
    stAng = math.degrees(theta1) * 60000
    swAng = math.degrees(dtheta) * 60000

    _emit_arc_to(path_el, rx, ry, stAng, swAng)
    return True


# ---------------------------------------------------------------------------
# Coordinate helpers
# ---------------------------------------------------------------------------

def _segments_to_emu(
    segments: list[PathSegment], scale: float
) -> list[PathSegment]:
    """Convert all segment coordinates from pixels to EMU."""
    factor = scale * EMU_PER_PX

    def c(v: float) -> float:
        return v * factor

    result: list[PathSegment] = []
    for seg in segments:
        if isinstance(seg, MoveToSeg):
            result.append(MoveToSeg(c(seg.x), c(seg.y)))
        elif isinstance(seg, LineToSeg):
            result.append(LineToSeg(c(seg.x), c(seg.y)))
        elif isinstance(seg, CubicBezierToSeg):
            result.append(CubicBezierToSeg(
                cp1=(c(seg.cp1[0]), c(seg.cp1[1])),
                cp2=(c(seg.cp2[0]), c(seg.cp2[1])),
                end=(c(seg.end[0]), c(seg.end[1])),
            ))
        elif isinstance(seg, QuadBezierToSeg):
            result.append(QuadBezierToSeg(
                cp=(c(seg.cp[0]), c(seg.cp[1])),
                end=(c(seg.end[0]), c(seg.end[1])),
            ))
        elif isinstance(seg, ArcSeg):
            result.append(ArcSeg(
                rx=c(seg.rx),
                ry=c(seg.ry),
                x_rotation=seg.x_rotation,
                large_arc=seg.large_arc,
                sweep=seg.sweep,
                x=c(seg.x),
                y=c(seg.y),
            ))
        elif isinstance(seg, CloseSeg):
            result.append(CloseSeg())
    return result


def _bounding_box(
    segments: list[PathSegment],
) -> tuple[Optional[float], Optional[float], Optional[float], Optional[float]]:
    """
    Compute the bounding box over all key points in segments.

    For bezier curves the convex hull of control points is used (conservative
    but always contains the curve).  For arcs a coarse sampling is used.
    """
    xs: list[float] = []
    ys: list[float] = []
    pen_x = pen_y = 0.0

    for seg in segments:
        if isinstance(seg, MoveToSeg):
            xs.append(seg.x)
            ys.append(seg.y)
            pen_x, pen_y = seg.x, seg.y

        elif isinstance(seg, LineToSeg):
            xs.append(seg.x)
            ys.append(seg.y)
            pen_x, pen_y = seg.x, seg.y

        elif isinstance(seg, CubicBezierToSeg):
            xs.extend([seg.cp1[0], seg.cp2[0], seg.end[0]])
            ys.extend([seg.cp1[1], seg.cp2[1], seg.end[1]])
            pen_x, pen_y = seg.end

        elif isinstance(seg, QuadBezierToSeg):
            xs.extend([seg.cp[0], seg.end[0]])
            ys.extend([seg.cp[1], seg.end[1]])
            pen_x, pen_y = seg.end

        elif isinstance(seg, ArcSeg):
            pts = svg_arc_to_lines(
                x1=pen_x, y1=pen_y,
                rx=seg.rx, ry=seg.ry,
                x_axis_rotation=seg.x_rotation,
                large_arc_flag=seg.large_arc,
                sweep_flag=seg.sweep,
                x2=seg.x, y2=seg.y,
                tolerance=float(EMU_PER_PX) * 10,  # coarse for bbox only
            )
            for px, py in pts:
                xs.append(px)
                ys.append(py)
            pen_x, pen_y = seg.x, seg.y

    if not xs:
        return None, None, None, None

    return min(xs), min(ys), max(xs), max(ys)
