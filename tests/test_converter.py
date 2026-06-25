"""Tests for the main SVG to PPTX converter."""

import pytest
from pathlib import Path
import tempfile
import os

from pptx import Presentation

from svg2pptx import svg_to_pptx, SVGConverter, Config


FIXTURES_DIR = Path(__file__).parent / "fixtures"


class TestSVGConverter:
    """Tests for SVGConverter class."""

    def test_convert_simple_svg_string(self):
        """Test converting a simple SVG string."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="red"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        assert prs is not None
        assert len(prs.slides) == 1
        # Should have at least one shape
        assert len(prs.slides[0].shapes) >= 1

    def test_convert_circle(self):
        """Test converting a circle."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <circle cx="50" cy="50" r="40" fill="blue"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        assert len(prs.slides[0].shapes) >= 1

    def test_convert_with_config(self):
        """Test conversion with custom config."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="0" y="0" width="50" height="50"/>
        </svg>'''
        
        config = Config(scale=2.0)
        converter = SVGConverter(config=config)
        prs = converter.convert_string(svg)
        
        assert prs is not None

    def test_convert_multiple_shapes(self):
        """Test converting multiple shapes."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">
            <rect x="10" y="10" width="50" height="50" fill="red"/>
            <circle cx="120" cy="35" r="25" fill="blue"/>
            <ellipse cx="170" cy="60" rx="20" ry="30" fill="green"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        # Should have 3 shapes
        assert len(prs.slides[0].shapes) >= 3


class TestConvertFile:
    """Tests for file-based conversion."""

    def test_convert_basic_shapes(self):
        """Test converting the basic shapes fixture."""
        svg_path = FIXTURES_DIR / "basic_shapes.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            svg_to_pptx(str(svg_path), pptx_path)
            
            # Verify the output file exists and is valid
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides) == 1
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    def test_convert_path_icon(self):
        """Test converting a path-based icon."""
        svg_path = FIXTURES_DIR / "path_icon.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            svg_to_pptx(str(svg_path), pptx_path)
            
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    def test_convert_grouped_svg(self):
        """Test converting SVG with groups."""
        svg_path = FIXTURES_DIR / "grouped.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            config = Config(preserve_groups=True)
            svg_to_pptx(str(svg_path), pptx_path, config=config)
            
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)


class TestAddToSlide:
    """Tests for adding SVG to existing slides."""

    def test_add_to_existing_slide(self):
        """Test adding SVG shapes to an existing presentation."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="0" y="0" width="50" height="50" fill="red"/>
        </svg>'''
        
        # Create a presentation with one slide
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add SVG shapes
        converter = SVGConverter()
        converter.add_string_to_slide(svg, slide)
        
        # Should have at least one shape
        assert len(slide.shapes) >= 1

    def test_add_multiple_svgs(self):
        """Test adding multiple SVGs to the same slide."""
        svg1 = '''<svg xmlns="http://www.w3.org/2000/svg" width="50" height="50">
            <rect x="0" y="0" width="50" height="50" fill="red"/>
        </svg>'''
        svg2 = '''<svg xmlns="http://www.w3.org/2000/svg" width="50" height="50">
            <circle cx="25" cy="25" r="25" fill="blue"/>
        </svg>'''
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        converter = SVGConverter()
        converter.add_string_to_slide(svg1, slide, x=0, y=0)
        converter.add_string_to_slide(svg2, slide, x=500000, y=0)
        
        # Should have at least 2 shapes
        assert len(slide.shapes) >= 2


class TestGradientOnRotatedShapes:
    """Regression tests for gradient direction on rotated shapes."""

    def _get_gradient_angle(self, prs) -> float:
        """Extract gradient angle from the first gradient-filled shape."""
        import lxml.etree as etree
        NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for shape in prs.slides[0].shapes:
            sp = shape._element
            for gradFill in sp.iter(f"{{{NS}}}gradFill"):
                lin = gradFill.find(f"{{{NS}}}lin")
                if lin is not None:
                    ang = int(lin.get("ang", "0"))
                    return ang / 60000.0  # convert from 1/60000 degree units
        return None

    def test_gradient_direction_unrotated_rect(self):
        """A horizontal gradient on a non-rotated rect stays horizontal (angle≈0°)."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg"
                      xmlns:xlink="http://www.w3.org/1999/xlink"
                      width="100" height="50">
          <defs>
            <linearGradient id="g1">
              <stop offset="0" style="stop-color:#000000;stop-opacity:1"/>
              <stop offset="1" style="stop-color:#ffffff;stop-opacity:1"/>
            </linearGradient>
            <linearGradient xlink:href="#g1" id="g2"
              x1="0" y1="0" x2="100" y2="0" gradientUnits="userSpaceOnUse"/>
          </defs>
          <rect x="0" y="0" width="100" height="50" fill="url(#g2)"/>
        </svg>'''
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        angle = self._get_gradient_angle(prs)
        assert angle is not None
        # horizontal gradient → angle ≈ 0°
        assert abs(angle % 360) < 1.0 or abs((angle % 360) - 360) < 1.0

    def test_gradient_direction_rotated90_rect(self):
        """A horizontal gradient on a rect rotated 90° becomes vertical (different from 0°)."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg"
                      xmlns:xlink="http://www.w3.org/1999/xlink"
                      width="50" height="100">
          <defs>
            <linearGradient id="g1">
              <stop offset="0" style="stop-color:#000000;stop-opacity:1"/>
              <stop offset="1" style="stop-color:#ffffff;stop-opacity:1"/>
            </linearGradient>
            <linearGradient xlink:href="#g1" id="g2"
              x1="0" y1="0" x2="100" y2="0" gradientUnits="userSpaceOnUse"/>
          </defs>
          <rect x="0" y="0" width="100" height="50" fill="url(#g2)"
                transform="rotate(90)"/>
        </svg>'''
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        angle = self._get_gradient_angle(prs)
        assert angle is not None
        # horizontal in local space + rotate(90) → vertical gradient, raw ang value = 90°
        assert abs((angle % 360) - 90.0) < 1.0

    def test_gradient_direction_rotated_minus90_rect(self):
        """A horizontal gradient on a rect rotated -90° becomes vertical (opposite to rotate(90))."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg"
                      xmlns:xlink="http://www.w3.org/1999/xlink"
                      width="50" height="100">
          <defs>
            <linearGradient id="g1">
              <stop offset="0" style="stop-color:#000000;stop-opacity:1"/>
              <stop offset="1" style="stop-color:#ffffff;stop-opacity:1"/>
            </linearGradient>
            <linearGradient xlink:href="#g1" id="g2"
              x1="0" y1="0" x2="100" y2="0" gradientUnits="userSpaceOnUse"/>
          </defs>
          <rect x="0" y="0" width="100" height="50" fill="url(#g2)"
                transform="rotate(-90)"/>
        </svg>'''
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        angle = self._get_gradient_angle(prs)
        assert angle is not None
        # horizontal in local space + rotate(-90) → vertical gradient, raw ang value = 270°
        assert abs((angle % 360) - 270.0) < 1.0


class TestRadialAndStrokeGradients:
    """Tests for radial gradients and gradients applied to strokes."""

    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _gradfills(self, prs):
        """Yield (parent_tag, gradFill_element) for every gradient in slide 0."""
        for shape in prs.slides[0].shapes:
            for gradFill in shape._element.iter(f"{{{self.NS}}}gradFill"):
                parent = gradFill.getparent().tag.split("}")[-1]
                yield parent, gradFill

    def test_radial_fill_emits_path_gradient(self):
        """A radialGradient fill becomes an OOXML <a:path path='circle'>."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
          <defs>
            <radialGradient id="rg">
              <stop offset="0" stop-color="#ff0000"/>
              <stop offset="1" stop-color="#0000ff"/>
            </radialGradient>
          </defs>
          <rect x="0" y="0" width="100" height="100" fill="url(#rg)"/>
        </svg>'''
        prs = SVGConverter().convert_string(svg)
        gradients = list(self._gradfills(prs))
        assert len(gradients) == 1
        parent, gradFill = gradients[0]
        assert parent == "spPr"  # a shape fill
        path = gradFill.find(f"{{{self.NS}}}path")
        assert path is not None and path.get("path") == "circle"
        assert gradFill.find(f"{{{self.NS}}}lin") is None
        stops = gradFill.findall(f"{{{self.NS}}}gsLst/{{{self.NS}}}gs")
        assert len(stops) == 2

    def test_radial_focus_shifted_by_center(self):
        """cx/cy shift the fillToRect focus for objectBoundingBox units."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
          <defs>
            <radialGradient id="rg" cx="0.3" cy="0.4">
              <stop offset="0" stop-color="#ff0000"/>
              <stop offset="1" stop-color="#0000ff"/>
            </radialGradient>
          </defs>
          <rect x="0" y="0" width="100" height="100" fill="url(#rg)"/>
        </svg>'''
        prs = SVGConverter().convert_string(svg)
        _, gradFill = next(self._gradfills(prs))
        ftr = gradFill.find(f"{{{self.NS}}}path/{{{self.NS}}}fillToRect")
        assert ftr is not None
        assert ftr.get("l") == "30000"
        assert ftr.get("t") == "40000"
        assert ftr.get("r") == "70000"
        assert ftr.get("b") == "60000"

    def test_radial_inherits_via_href(self):
        """A radialGradient inherits stops and geometry through xlink:href."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg"
                      xmlns:xlink="http://www.w3.org/1999/xlink"
                      width="100" height="100">
          <defs>
            <radialGradient id="base" cx="0.25" cy="0.75">
              <stop offset="0" stop-color="#00ff00"/>
              <stop offset="1" stop-color="#000000"/>
            </radialGradient>
            <radialGradient id="child" xlink:href="#base"/>
          </defs>
          <rect x="0" y="0" width="100" height="100" fill="url(#child)"/>
        </svg>'''
        prs = SVGConverter().convert_string(svg)
        _, gradFill = next(self._gradfills(prs))
        stops = gradFill.findall(f"{{{self.NS}}}gsLst/{{{self.NS}}}gs")
        assert len(stops) == 2
        ftr = gradFill.find(f"{{{self.NS}}}path/{{{self.NS}}}fillToRect")
        assert ftr.get("l") == "25000"  # cx=0.25 inherited
        assert ftr.get("t") == "75000"  # cy=0.75 inherited

    def test_gradient_on_stroke_emits_line_gradient(self):
        """A gradient stroke produces a <a:gradFill> inside the line's <a:ln>."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg"
                      xmlns:xlink="http://www.w3.org/1999/xlink"
                      width="100" height="100">
          <defs>
            <linearGradient id="lg">
              <stop offset="0" stop-color="#000000"/>
              <stop offset="1" stop-color="#ffffff"/>
            </linearGradient>
          </defs>
          <rect x="10" y="10" width="80" height="80"
                fill="none" stroke="url(#lg)" stroke-width="4"/>
        </svg>'''
        prs = SVGConverter().convert_string(svg)
        gradients = list(self._gradfills(prs))
        assert len(gradients) == 1
        parent, gradFill = gradients[0]
        assert parent == "ln"  # the gradient lives on the line, not the fill
        stops = gradFill.findall(f"{{{self.NS}}}gsLst/{{{self.NS}}}gs")
        assert len(stops) == 2

    def test_gradient_stroke_on_line_connector(self):
        """A gradient stroke on a <line> renders as a native connector gradient.

        Regression: <line> uses the connector path (not apply_style), so a
        gradient stroke must still produce an <a:gradFill> inside the line's
        <a:ln> rather than silently rendering with no color.
        """
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
          <defs>
            <linearGradient id="lg">
              <stop offset="0" stop-color="#000000"/>
              <stop offset="1" stop-color="#ffffff"/>
            </linearGradient>
          </defs>
          <line x1="0" y1="0" x2="100" y2="100"
                stroke="url(#lg)" stroke-width="3"/>
        </svg>'''
        prs = SVGConverter().convert_string(svg)
        gradients = list(self._gradfills(prs))
        assert len(gradients) == 1
        parent, gradFill = gradients[0]
        assert parent == "ln"  # gradient lives on the connector's line
        stops = gradFill.findall(f"{{{self.NS}}}gsLst/{{{self.NS}}}gs")
        assert len(stops) == 2

    def test_radial_pptx_roundtrips(self):
        """A presentation with radial fill + stroke reloads without error."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
          <defs>
            <radialGradient id="rg">
              <stop offset="0" stop-color="#ff0000"/>
              <stop offset="1" stop-color="#0000ff" stop-opacity="0.5"/>
            </radialGradient>
          </defs>
          <rect x="0" y="0" width="100" height="100" fill="url(#rg)"/>
          <circle cx="50" cy="50" r="20" fill="none"
                  stroke="url(#rg)" stroke-width="5"/>
        </svg>'''
        prs = SVGConverter().convert_string(svg)
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            reopened = Presentation(path)
            assert len(reopened.slides) == 1
        finally:
            os.unlink(path)


class TestConfig:
    """Tests for configuration options."""

    def test_default_config(self):
        config = Config()
        assert config.scale == 1.0
        assert config.preserve_groups is False
        assert config.flatten_groups is True
        assert config.disable_shadows is True

    def test_custom_config(self):
        config = Config(
            scale=2.0,
            curve_tolerance=0.5,
            preserve_groups=False,
        )
        assert config.scale == 2.0
        assert config.curve_tolerance == 0.5
        assert config.preserve_groups is False
